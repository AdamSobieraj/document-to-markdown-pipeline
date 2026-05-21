import io
import logging
import os
from typing import List, Union
from pptx import Presentation
from langchain_core.documents import Document
from .base_parser import BaseDocumentParser

logger = logging.getLogger(__name__)


class PptxParser(BaseDocumentParser):
    """
    Parser dla plików PowerPoint (.pptx).
    Każdy slajd to oddzielny Document. Zamienia tekst i tabele na Markdown.
    """

    def parse(self, file_source: Union[str, io.BytesIO, bytes], **kwargs) -> List[Document]:
        documents = []
        source_name = os.path.basename(file_source) if isinstance(file_source, str) else "strumień pamięci S3"

        try:
            # python-pptx wymaga obiektu plikopodobnego (str lub BytesIO)
            # Jeśli dostaniemy surowe bajty (bytes), owijamy je w BytesIO
            if isinstance(file_source, bytes):
                file_source = io.BytesIO(file_source)

            prs = Presentation(file_source)

            # Iteracja po slajdach (każdy slajd to 1 strona/Document)
            for i, slide in enumerate(prs.slides):
                md_lines = []

                # Oznaczamy numer slajdu jako nagłówek na potrzeby czytelności Markdown
                md_lines.append(f"## Slajd {i + 1}")

                for shape in slide.shapes:
                    # 1. Wyciąganie tekstu z pól tekstowych
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                # Prosta heurystyka dla list (jeśli akapit ma poziom > 0, traktuj jak listę)
                                if paragraph.level > 0:
                                    md_lines.append(f"- {text}")
                                else:
                                    md_lines.append(text)

                    # 2. Wyciąganie tabel do formatu tabel Markdown
                    elif shape.has_table:
                        table = shape.table
                        for row_idx, row in enumerate(table.rows):
                            row_data = []
                            for cell in row.cells:
                                # Pobieramy tekst z komórki, usuwamy entery i czyścimy
                                cell_text = cell.text_frame.text.replace('\n', ' ').strip()
                                row_data.append(cell_text)

                            # Budowa wiersza: | Wartość 1 | Wartość 2 |
                            md_row = "| " + " | ".join(row_data) + " |"
                            md_lines.append(md_row)

                            # Budowa separatora tabeli pod pierwszym wierszem (nagłówkiem)
                            if row_idx == 0:
                                separator = "|-" + "-|-".join(
                                    ["-" * len(c) if len(c) > 0 else "-" for c in row_data]) + "-|"
                                md_lines.append(separator)

                # Złącz linie dla całego slajdu podwójnym enterem dla zachowania bloków Markdown
                if len(md_lines) > 1:  # Ignoruj puste slajdy (mają tylko nagłówek "## Slajd X")
                    slide_text = "\n\n".join(md_lines).strip()

                    doc = Document(
                        page_content=slide_text,
                        metadata={"page_number": i + 1}
                    )
                    documents.append(doc)

            return documents

        except Exception as e:
            logger.error(f"Błąd przetwarzania PPTX ({source_name}): {e}")
            return []